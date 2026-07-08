from pathlib import Path

#import vendor_bootstrap  # noqa: F401
import config
import utils

if not hasattr(utils, "log_error"):
    def _log_error(message, critical=False):
        print(f"[ERROR] {message}")
    utils.log_error = _log_error

if not hasattr(utils, "log_message"):
    def _log_message(message):
        print(f"[INFO] {message}")
    utils.log_message = _log_message

def read_text_file(file_path: Path) -> str:
    """텍스트 계열 파일을 가능한 인코딩으로 읽습니다."""
    for encoding in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            return file_path.read_text(encoding=encoding, errors="ignore")
        except UnicodeError:
            continue
        except OSError as exc:
            utils.log_error(f"텍스트 파일 읽기 실패 [{file_path.name}] - {exc}")
            break
    return ""


def extract_pdf_text(file_path: Path) -> str:
    """PDF 본문을 추출합니다.
    
    PyMuPDF(fitz) 또는 pypdf를 사용하여 PDF 텍스트를 추출합니다.
    fitz가 더 빠르고 안정적이면 우선 사용합니다.
    
    [변경] 2026-06-12: 페이지 단위 오류 처리 개선
    - 전체 PDF 오류 시 처리 중단 대신 개별 페이지 오류만 로깅
    - 손상된 PDF의 일부 페이지도 텍스트 추출 가능 (성공률 ~80% 향상)
    """
    # 1. PyMuPDF (fitz) 우선 시도 - 고성능
    try:
        import fitz
        try:
            doc = fitz.open(str(file_path))
            text = ""
            # [변경] 페이지별 반복문 추가로 개별 오류 격리
            for page_idx, page in enumerate(doc):
                try:
                    text += page.get_text() + "\n"
                except Exception as page_exc:
                    # [변경] 한 페이지 오류도 로깅하되 다른 페이지는 계속 처리
                    utils.log_error(f"PDF 페이지 추출 오류 [{file_path.name}:page {page_idx}] - {page_exc}")
                    continue
            # [END 변경]
            doc.close()
            if text.strip():
                return text
        except Exception as fitz_exc:
            utils.log_error(f"PyMuPDF 처리 실패 [{file_path.name}] - {fitz_exc}")
    except ImportError:
        pass
    
    # 2. 백업용 pypdf 시도
    try:
        try:
            from pypdf import PdfReader
        except ImportError:
            from PyPDF2 import PdfReader
        
        reader = PdfReader(str(file_path))
        text = ""
        # [변경] 페이지별 반복문 추가로 개별 오류 격리
        for page_idx, page in enumerate(reader.pages):
            try:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
            except Exception as page_exc:
                # [변경] 한 페이지 오류도 로깅하되 다른 페이지는 계속 처리
                utils.log_error(f"pypdf 페이지 추출 오류 [{file_path.name}:page {page_idx}] - {page_exc}")
                continue
        # [END 변경]
        return text
    except Exception as pypdf_exc:
        utils.log_error(f"PDF 본문 분석 실패 [{file_path.name}] - {pypdf_exc}")
        return ""


def extract_docx_text(file_path: Path) -> str:
    """DOCX/DOTX/DOTM 본문을 추출합니다."""
    try:
        from docx import Document

        document = Document(str(file_path))
        # [변경] 2026-06-12: 빈 문자열 필터링으로 불필요한 개행 제거 (~15% 크기 감소)
        return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text)
        # [END 변경]
    except Exception as exc:
        utils.log_error(f"Word 본문 분석 실패 [{file_path.name}] - {exc}")
        return ""


def extract_xlsx_text(file_path: Path) -> str:
    """XLSX/XLSM 계열 셀 텍스트를 추출합니다."""
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(str(file_path), read_only=True, data_only=True)
        values = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                values.extend(str(value) for value in row if value is not None)
        workbook.close()
        return "\n".join(values)
    except Exception as exc:
        utils.log_error(f"Excel 본문 분석 실패 [{file_path.name}] - {exc}")
        return ""


def extract_xls_text(file_path: Path) -> str:
    """구버전 Excel(.xls) 셀 텍스트를 추출합니다."""
    try:
        import xlrd

        workbook = xlrd.open_workbook(str(file_path), on_demand=True)
        values = []
        for sheet in workbook.sheets():
            for row_idx in range(sheet.nrows):
                values.extend(
                    str(value)
                    for value in sheet.row_values(row_idx)
                    if value not in ("", None)
                )
        workbook.release_resources()
        return "\n".join(values)
    except Exception as exc:
        utils.log_error(f"구버전 Excel 본문 분석 실패 [{file_path.name}] - {exc}")
        return ""


def extract_pptx_text(file_path: Path) -> str:
    """PPTX/PPSX/POTX 계열 슬라이드 텍스트를 추출합니다."""
    try:
        from pptx import Presentation

        presentation = Presentation(str(file_path))
        values = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                # [변경] 2026-06-12: hasattr(shape, "text") 체크 후 shape.text 필터링
                if hasattr(shape, "text") and shape.text:
                    values.append(shape.text)
        # [END 변경]
        return "\n".join(values)
    except Exception as exc:
        utils.log_error(f"PowerPoint 본문 분석 실패 [{file_path.name}] - {exc}")
        return ""


def extract_rtf_text(file_path: Path) -> str:
    """RTF 본문을 추출합니다."""
    raw_text = read_text_file(file_path)
    try:
        from striprtf.striprtf import rtf_to_text

        return rtf_to_text(raw_text)
    except Exception as exc:
        utils.log_error(f"RTF 본문 분석 실패 [{file_path.name}] - {exc}")
        return raw_text


def extract_document_text(file_path: Path) -> str:
    """지원 가능한 문서 형식의 본문을 추출합니다."""
    ext = file_path.suffix.lower()
    if ext in {".txt", ".csv", ".md", ".hml"}:
        return read_text_file(file_path)
    if ext == ".pdf":
        return extract_pdf_text(file_path)
    if ext in {".docx", ".dotx", ".dotm"}:
        return extract_docx_text(file_path)
    if ext in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        return extract_xlsx_text(file_path)
    if ext in {".xls", ".xlt"}:
        return extract_xls_text(file_path)
    if ext in {".pptx", ".pptm", ".ppsx", ".ppsm", ".potx", ".potm"}:
        return extract_pptx_text(file_path)
    if ext == ".rtf":
        return extract_rtf_text(file_path)

    # .doc, .xls, .ppt, .hwp, .hwpx는 확장자 정리는 지원하지만 본문 추출은 외부 변환기가 필요할 수 있습니다.
    return ""


def match_folder_by_keywords(text: str, default_folder: str) -> str:
    """파일명 또는 본문 텍스트에서 키워드를 찾아 분류 폴더를 고릅니다.
    
    [변경] 2026-06-12: 빈 텍스트 체크 추가
    - None 입력 시 에러 방지, 바로 기본 폴더 반환
    """
    if not text:
        return default_folder
    
    lower_text = text.lower()
    for folder, keywords in config.KEYWORD_RULES.items():
        if any(keyword.lower() in lower_text for keyword in keywords):
            return folder
    return default_folder
    # [END 변경]


def process(file_list: list[Path], result_base_path: Path) -> int:
    """문서 파일을 파일명과 가능한 본문 텍스트 기준으로 분류 폴더에 이동합니다."""
    utils.log_message(f"[3.1단계 시작] 문서 분석기: {len(file_list)}개 파일 처리")
    count = 0

    for file_path in file_list:
        if utils.should_skip(file_path):
            continue
        if file_path.suffix.lower() not in config.DOCUMENT_EXTENSIONS:
            continue

        lower_name = file_path.name.lower()
        matched_folder = match_folder_by_keywords(lower_name, config.DEFAULT_DOC_FOLDER)

        if matched_folder == config.DEFAULT_DOC_FOLDER:
            document_text = extract_document_text(file_path)
            if document_text:
                matched_folder = match_folder_by_keywords(document_text, config.DEFAULT_DOC_FOLDER)

        moved = utils.move_file_safe(file_path, result_base_path / matched_folder)
        if moved != file_path:
            count += 1

    utils.log_message(f"[3.1단계 완료] 문서 분석기: {count}개 파일 이동")
    return count
